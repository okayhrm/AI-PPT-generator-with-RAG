# ppt_gen.py — Chat-aware pipeline (Gemini + Pinecone v5 + local embeddings) • Py3.13
from __future__ import annotations
import os, re, json, time, datetime, tempfile
from typing import List, Tuple, TypedDict, Optional

from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt

import google.generativeai as genai
from langgraph.graph import StateGraph, END

load_dotenv()
# Accept either key name
_G_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY") or ""

TEMPLATES_DIR    = os.getenv("TEMPLATES_DIR", "./templates")
DEFAULT_TEMPLATE = "template1.pptx"
DEPTH_DEFAULT    = "detailed"
DEFAULT_SLIDES   = 7
MAX_SLIDES       = 20

PINECONE_INDEX  = os.getenv("PINECONE_INDEX", "ppt-rag")

# ---------- small helpers ----------
def _n_slides(prompt: str) -> int:
    m = re.search(r"(\d+)\s*[- ]*slide", (prompt or "").lower())
    n = int(m.group(1)) if m else DEFAULT_SLIDES
    return max(1, min(n, MAX_SLIDES))

def _depth_to_rules(depth: str) -> tuple[int, int]:
    depth = (depth or DEPTH_DEFAULT).lower()
    if depth == "light":     return (3, 10)
    if depth == "extensive": return (7, 18)
    return (4, 14)

def _get_gemini_model(model_name: str = "gemini-2.5-flash"):
    if not _G_KEY:
        raise RuntimeError("GEMINI_API_KEY / GOOGLE_API_KEY not set")
    genai.configure(api_key=_G_KEY)
    return genai.GenerativeModel(model_name)

def _safe_json_extract(text: str) -> dict:
    i = text.find("{"); j = text.rfind("}")
    if i == -1 or j == -1 or j <= i:
        raise ValueError("Model did not return JSON.")
    return json.loads(text[i:j+1])

# ---------- PPT builder ----------
def build_ppt(slides: List[Tuple[str, List[str]]], template_filename: str) -> str:
    path = os.path.join(TEMPLATES_DIR, template_filename or DEFAULT_TEMPLATE)
    prs = Presentation(path)

    # find title+body layout
    layout = None
    for l in prs.slide_layouts:
        has_title = any(getattr(ph.placeholder_format, "type", None) == 1 for ph in l.placeholders)
        has_body  = any(getattr(ph.placeholder_format, "type", None) == 2 for ph in l.placeholders)
        if has_title and has_body:
            layout = l; break
    if layout is None:
        layout = prs.slide_layouts[1]

    # remove lone cover if desired
    if len(prs.slides) == 1 and len(prs.slides._sldIdLst) == 1:
        rid = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[0]

    for title, bullets in slides:
        sl = prs.slides.add_slide(layout)
        if sl.shapes.title:
            sl.shapes.title.text = title
        body = None
        for sh in sl.placeholders:
            if getattr(sh.placeholder_format, "type", None) == 2 and sh.has_text_frame:
                body = sh; break
        if body:
            tf = body.text_frame; tf.clear()
            for b in bullets:
                p = tf.add_paragraph(); p.text = b; p.level = 0; p.font.size = Pt(18)

    out = os.path.join(tempfile.gettempdir(), f"deck_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx")
    prs.save(out)
    return out

# ---------- Pinecone RAG (v5 safe) ----------
def pinecone_query(text_query: str, namespace: str, top_k: int = 6) -> List[str]:
    from pinecone import Pinecone
    from embedder import get_embedder

    pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
    index = pc.Index(os.getenv("PINECONE_INDEX", PINECONE_INDEX))

    emb = get_embedder()
    qvec = emb.embed_query(text_query)
    res = index.query(namespace=namespace, vector=qvec, top_k=top_k, include_metadata=True)

    hits = []
    for m in (res.matches or []):
        md = m.get("metadata") or {}
        hits.append(md.get("content", ""))
    return [h for h in hits if h]

# ---------- Chat folding ----------
def history_to_instructions(history: list, max_chars: int = 1200) -> str:
    if not history: return ""
    buf = []
    for msg in reversed(history[-20:]):
        role = msg.get("role")
        txt  = (msg.get("content") or "").strip()
        if role == "user" and txt:
            buf.append(f"- {txt}")
    block = "\n".join(buf[:10])
    return block[:max_chars]

# ---------- Slide generation ----------
def generate_slides(topic: str, n: int, depth: str = DEPTH_DEFAULT, ctx: str = "", history: list | None = None) -> List[Tuple[str, List[str]]]:
    bullets_per, max_words = _depth_to_rules(depth)
    user_instructions = history_to_instructions(history or [])

    sys_prompt = (
        f"You are an expert presentation creator.\n"
        f"Generate exactly {n} slides in JSON only:\n"
        f'{{"slides":[{{"title":"Slide 1 Title","content":["point 1","point 2","point 3"]}}]}}\n'
        f"Rules:\n"
        f"- Each slide has {bullets_per} bullet points (±1 to fit n overall).\n"
        f"- Each bullet ≤ {max_words} words.\n"
        f"- Use crisp, factual language.\n"
        f"- Apply the user's chat instructions if provided (style, tone, structure).\n"
        f"- If context is provided, ground facts to it and avoid hallucinations.\n"
        f"- Output JSON only (no markdown).\n"
    )

    payload = topic
    if user_instructions:
        payload += f"\n\nUser Instructions (apply faithfully):\n{user_instructions}"
    if ctx:
        payload += f"\n\nContext (verbatim excerpts; prefer these facts):\n{ctx[:6000]}"

    model = _get_gemini_model("gemini-2.5-flash")
    for _ in range(3):
        rsp = model.generate_content(sys_prompt + "\n\n" + payload)
        raw = rsp.text or ""
        try:
            js = _safe_json_extract(raw)
            slides = js.get("slides", [])
            out: List[Tuple[str, List[str]]] = []
            for s in slides:
                title = (s.get("title") or "").strip()[:120]
                points = [(p or "").strip() for p in (s.get("content") or []) if (p or "").strip()]
                trim = []
                for p in points[: bullets_per + 1]:
                    w = p.split()
                    if len(w) > max_words: w = w[:max_words]
                    trim.append(" ".join(w))
                out.append((title, trim))
            if out:
                return out[:MAX_SLIDES]
        except Exception:
            time.sleep(0.5)
    raise RuntimeError("Gemini failed to return valid JSON 3×")

# ---------- Graph ----------
class State(TypedDict, total=False):
    prompt: str
    template: str
    depth: str
    route: str
    context: str
    slides: List[Tuple[str, List[str]]]
    output_path: str
    namespace: str
    force_rag: bool
    history: list

class Config(TypedDict, total=False):
    model_name: str
    selected_template: str
    ui_mode: bool
    force_rag: bool

def router_node(state: State, config: Optional[Config] = None):
    state["route"] = "rag" if state.get("force_rag") else "none"
    return state

def rag_node(state: State, config: Optional[Config] = None):
    ns = state.get("namespace") or state.get("project_id") or "default"
    chunks = pinecone_query(state["prompt"], namespace=ns, top_k=6)
    return {"context": "\n\n".join(chunks)}

def slide_gen_node(state: State, config: Optional[Config] = None):
    ctx = state.get("context", "")
    slides = generate_slides(
        topic=state["prompt"],
        n=_n_slides(state["prompt"]),
        depth=state.get("depth") or DEPTH_DEFAULT,
        ctx=ctx,
        history=state.get("history") or [],
    )
    return {"slides": slides}

def confirm_node(state: State, config: Optional[Config] = None):
    return {}

def build_ppt_node(state: State, config: Optional[Config] = None):
    cfg = (config or {})
    template = (cfg.get("selected_template") or state.get("template") or DEFAULT_TEMPLATE)
    path = build_ppt(state["slides"], template)
    return {"output_path": path}

g = StateGraph(State, config_schema=Config)
g.add_node("router_node", router_node)
g.add_node("rag_node",    rag_node)
g.add_node("slide_gen",   slide_gen_node)
g.add_node("confirm",     confirm_node)
g.add_node("build_ppt",   build_ppt_node)

g.set_entry_point("router_node")
g.add_conditional_edges("router_node", lambda st: st.get("route"), {"rag": "rag_node", "none": "slide_gen"})
g.add_edge("rag_node",  "slide_gen")
g.add_edge("slide_gen", "confirm")
g.add_edge("confirm",   "build_ppt")
g.add_edge("build_ppt", END)

graph = g.compile(name="ppt_graph_chat")
